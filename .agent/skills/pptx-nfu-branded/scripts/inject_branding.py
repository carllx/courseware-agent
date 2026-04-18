#!/usr/bin/env python3
"""
NFU (广州南方学院) PPT 品牌注入脚本

工作模式："固定环节 + 内容合并"
  - 固定环节：封面、回顾、目录、引用、作业、结尾（由本脚本生成）
  - 内容主体：由用户已有的课程 PPT 提供
  - 可选环节：章节分隔页

最终输出：[封面] → [回顾] → [目录] → [原有内容...] → [引用] → [作业] → [结尾]

使用方式：
  python inject_branding.py --input content.pptx --output branded.pptx \\
      --course-name "交互设计" --teacher "张三" ...

依赖：python-pptx, pyyaml, lxml
"""

import argparse
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    print("错误：缺少依赖。请执行：pip install python-pptx lxml")
    sys.exit(1)

try:
    import yaml
except ImportError:
    print("错误：缺少 PyYAML。请执行：pip install pyyaml")
    sys.exit(1)

# ====================================================================
# 路径解析（确保跨项目可迁移—全部相对于脚本自身位置）
# ====================================================================
SKILL_DIR = Path(__file__).resolve().parent.parent
RESOURCES_DIR = SKILL_DIR / "resources"
THEME_PATH = RESOURCES_DIR / "nfu_theme.yaml"


def load_theme() -> dict:
    """加载 NFU 设计 Token"""
    with open(THEME_PATH, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def resolve_media_dir(aspect: str) -> Path:
    """根据比例选择媒体素材目录"""
    return RESOURCES_DIR / "media" / ("16-9" if aspect == "16:9" else "4-3")


def detect_aspect(prs: Presentation) -> str:
    """探测幻灯片宽高比"""
    return "16:9" if (prs.slide_width / prs.slide_height) > 1.5 else "4:3"


# ====================================================================
# OOXML 低级工具
# ====================================================================

def _hex_rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def _find_blank_layout(prs: Presentation):
    """定位空白布局（优先 layoutType='blank' / 第 7 个 / 最后一个）"""
    for layout in prs.slide_layouts:
        el = layout.element
        if el.get("type") == "blank":
            return layout
    if len(prs.slide_layouts) > 6:
        return prs.slide_layouts[6]
    return prs.slide_layouts[-1]


def _move_slide(prs: Presentation, from_idx: int, to_idx: int):
    """将幻灯片从 from_idx 移动到 to_idx"""
    sldIdLst = prs.element.find(qn("p:sldIdLst"))
    entries = list(sldIdLst)
    target = entries[from_idx]
    sldIdLst.remove(target)
    if to_idx == 0:
        sldIdLst.insert(0, target)
    else:
        list(sldIdLst)[to_idx - 1].addnext(target)


def _set_ea_font(run, font_name: str):
    """为 run 设置东亚字体"""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_name)


def _set_line_spacing(paragraph, pct: int):
    """设置段落行距百分比"""
    pPr = paragraph._p.get_or_add_pPr()
    lnSpc = pPr.find(qn("a:lnSpc"))
    if lnSpc is None:
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    for child in list(lnSpc):
        lnSpc.remove(child)
    spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
    spcPct.set("val", str(pct * 1000))


# ====================================================================
# 通用文本构建器
# ====================================================================

def _add_fullscreen_bg(slide, image_path: str, w: int, h: int):
    """添加全屏背景图片"""
    slide.shapes.add_picture(image_path, 0, 0, w, h)


def _add_logo(slide, logo_path: str, cfg: dict):
    """在 layout 配置指定位置放置 Logo"""
    lc = cfg["logo"]
    slide.shapes.add_picture(logo_path, lc["x"], lc["y"], lc["cx"], lc["cy"])


def _add_textbox(slide, pos: dict, lines: list, theme: dict,
                 first_line_sz_key: str = "cover_title_sz",
                 body_sz_key: str = "meta_sz",
                 color_hex: str = "FFFFFF",
                 bold: bool = False,
                 spacing_pct: int = 160):
    """通用文本框构建器"""
    typo = theme["typography"]
    cn_font = typo["cn_primary"]

    txBox = slide.shapes.add_textbox(pos["x"], pos["y"], pos["cx"], pos["cy"])
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT

        run = p.add_run()
        run.text = text

        # 字号：首行用大标题，其余用正文
        sz_key = first_line_sz_key if i == 0 else body_sz_key
        if sz_key in typo:
            run.font.size = Pt(typo[sz_key] / 100)

        run.font.name = cn_font
        run.font.bold = bold
        run.font.color.rgb = _hex_rgb(color_hex)
        _set_ea_font(run, cn_font)
        _set_line_spacing(p, spacing_pct)

    return txBox


# ====================================================================
# 固定环节页生成函数
# ====================================================================

def gen_cover(prs, theme, layout, media, info, *, position=0):
    """封面页：深灰背景 + 课程名/编号/教师/学期 + Logo"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))
    w, h = layout["slide_width"], layout["slide_height"]
    _add_fullscreen_bg(slide, str(media / theme["media"]["cover_bg"]["file"]), w, h)
    _add_logo(slide, str(media / theme["media"]["logo_combo"]["file"]), layout)

    lines = [
        info.get("course_name", "课程名称"),
        f"课程编号：{info.get('course_code', '')}",
        "",
        f"主讲教师：{info.get('teacher', '')}",
        f"开课单位：{theme['brand']['name_cn']}",
        f"授课学期：{info.get('semester', '')}",
        f"授课时间：{info.get('date', '')}",
    ]
    _add_textbox(slide, layout["cover_text"], lines, theme,
                 color_hex=theme["palette"]["text_on_dark"],
                 spacing_pct=160)

    if position == 0:
        _move_slide(prs, len(prs.slides) - 1, 0)
    return slide


def gen_review(prs, theme, layout, media, info, *,
               title="上节课知识点回顾", points=None):
    """回顾页：浅灰背景 + 知识点列表"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))
    w, h = layout["slide_width"], layout["slide_height"]
    _add_fullscreen_bg(slide, str(media / theme["media"]["content_bg"]["file"]), w, h)

    # 标题
    _add_textbox(slide, layout["review_title"], [title], theme,
                 first_line_sz_key="page_title_sz",
                 body_sz_key="page_title_sz",
                 color_hex=theme["palette"]["title_muted"],
                 spacing_pct=100)

    # 知识点列表
    if points:
        pt_lines = [f"知识点{i}：{p}" for i, p in enumerate(points, 1)]
        _add_textbox(slide, layout["review_content"], pt_lines, theme,
                     first_line_sz_key="section_title_sz",
                     body_sz_key="body_sz",
                     color_hex=theme["palette"]["text_on_light"],
                     spacing_pct=115)
    return slide


def gen_toc(prs, theme, layout, media, info, *, sections=None):
    """目录页：浅灰背景 + 章/节概述"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))
    w, h = layout["slide_width"], layout["slide_height"]
    _add_fullscreen_bg(slide, str(media / theme["media"]["content_bg"]["file"]), w, h)

    _add_textbox(slide, layout["review_title"], ["本节课主要内容概述"], theme,
                 first_line_sz_key="page_title_sz",
                 body_sz_key="page_title_sz",
                 color_hex=theme["palette"]["title_muted"],
                 spacing_pct=100)

    if sections:
        _add_textbox(slide, layout["review_content"], sections, theme,
                     first_line_sz_key="body_sz",
                     body_sz_key="body_sz",
                     color_hex=theme["palette"]["text_on_light"],
                     spacing_pct=150)
    return slide


def gen_section_header(prs, theme, layout, media, *,
                       chapter="", section_title=""):
    """章节分隔页（可选）：白色简洁"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))

    lines = []
    if chapter:
        lines.append(chapter)
    if section_title:
        lines.append(section_title)

    _add_textbox(slide, layout["section_meta"], lines, theme,
                 first_line_sz_key="small_sz",
                 body_sz_key="small_sz",
                 color_hex=theme["palette"]["text_on_light"],
                 spacing_pct=150)
    return slide


def gen_reference(prs, theme, layout, media, info, *,
                  references=None):
    """引用/拓展页：浅灰背景 + 参考列表"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))
    w, h = layout["slide_width"], layout["slide_height"]
    _add_fullscreen_bg(slide, str(media / theme["media"]["content_bg"]["file"]), w, h)

    _add_textbox(slide, layout["review_title"], ["本节课授课内容引用"], theme,
                 first_line_sz_key="page_title_sz",
                 body_sz_key="page_title_sz",
                 color_hex=theme["palette"]["title_muted"],
                 spacing_pct=100)

    if references:
        _add_textbox(slide, layout["review_content"], references, theme,
                     first_line_sz_key="body_sz",
                     body_sz_key="body_sz",
                     color_hex=theme["palette"]["text_on_light"],
                     spacing_pct=115)
    return slide


def gen_assignment(prs, theme, layout, media, info, *,
                   assignment_text=None):
    """作业页：浅灰背景 + 作业要求"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))
    w, h = layout["slide_width"], layout["slide_height"]
    _add_fullscreen_bg(slide, str(media / theme["media"]["content_bg"]["file"]), w, h)

    _add_textbox(slide, layout["review_title"], ["作业要求"], theme,
                 first_line_sz_key="page_title_sz",
                 body_sz_key="page_title_sz",
                 color_hex=theme["palette"]["title_muted"],
                 spacing_pct=100)

    if assignment_text:
        _add_textbox(slide, layout["review_content"], assignment_text, theme,
                     first_line_sz_key="body_sz",
                     body_sz_key="body_sz",
                     color_hex=theme["palette"]["text_on_light"],
                     bold=True,
                     spacing_pct=115)
    return slide


def gen_ending(prs, theme, layout, media, info):
    """结尾页：深灰背景 + 联系方式 + Logo"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))
    w, h = layout["slide_width"], layout["slide_height"]
    _add_fullscreen_bg(slide, str(media / theme["media"]["cover_bg"]["file"]), w, h)
    _add_logo(slide, str(media / theme["media"]["logo_combo"]["file"]), layout)

    lines = [
        info.get("course_name", "课程名称"),
        f"课程编号：{info.get('course_code', '')}",
        "",
        f"主讲教师：{info.get('teacher', '')}",
        f"办公邮箱：{info.get('email', '')}",
        f"办公地点：{info.get('office', '')}",
    ]
    _add_textbox(slide, layout["cover_text"], lines, theme,
                 color_hex=theme["palette"]["text_on_dark"],
                 bold=True,
                 spacing_pct=160)
    return slide


# ====================================================================
# 主题色注入
# ====================================================================

def inject_theme_colors(prs: Presentation, theme: dict):
    """替换 PPT 主题色板为 NFU 标准

    通过 slide_master → theme 关系链找到 theme part，
    再用 blob（原始 XML bytes）解析修改后回写。
    兼容 python-pptx 1.0.x 中通用 Part 无 element 属性的情况。
    """
    pal = theme["palette"]
    color_map = {
        "accent1": pal["accent1"], "accent2": pal["accent2"],
        "accent3": pal["accent3"], "accent4": pal["accent4"],
        "accent5": pal["accent5"], "accent6": pal["accent6"],
        "dk2": pal["dk2"], "lt2": pal["lt2"],
        "hlink": pal["hyperlink"], "folHlink": pal["followed_link"],
    }

    # 通过 slide master 的 rels 定位 theme part
    theme_part = None
    for sm in prs.slide_masters:
        for rId, rel in sm.part.rels.items():
            if "theme" in str(rel.target_ref):
                theme_part = rel.target_part
                break
        if theme_part:
            break

    if theme_part is None:
        raise ValueError("未找到 theme part")

    # 从 blob 读取 XML → 修改 → 回写
    root = etree.fromstring(theme_part.blob)
    cs = root.find(".//" + qn("a:clrScheme"))
    if cs is None:
        raise ValueError("theme XML 中未找到 clrScheme")

    for tag, hx in color_map.items():
        elem = cs.find(qn(f"a:{tag}"))
        if elem is not None:
            for ch in list(elem):
                elem.remove(ch)
            srgb = etree.SubElement(elem, qn("a:srgbClr"))
            srgb.set("val", hx)

    # 回写修改后的 XML
    theme_part._blob = etree.tostring(root, xml_declaration=True,
                                       encoding="UTF-8", standalone=True)


# ====================================================================
# course.yaml 读取
# ====================================================================

def load_course_yaml(path: str, week: int = None) -> dict:
    """从 course.yaml 提取课程信息

    支持嵌套结构：course.name, teacher.name 等。
    若提供 week，还会从 calendar 中提取该周的教学内容。
    """
    with open(path, "r", encoding="utf-8") as f:
        data = yaml.safe_load(f)

    # 兼容嵌套 / 扁平两种结构
    course = data.get("course", data.get("meta", data))
    teacher_obj = data.get("teacher", {})
    if isinstance(teacher_obj, str):
        teacher_name = teacher_obj
        teacher_office = ""
    else:
        teacher_name = teacher_obj.get("name", "")
        teacher_office = teacher_obj.get("office", "")

    # 学期显示格式化
    raw_sem = course.get("semester", "")
    if raw_sem and "-" in str(raw_sem):
        parts = str(raw_sem).split("-")
        if len(parts) == 3:
            semester_display = f"{parts[0]}-{parts[1]}学年第{'一' if parts[2] == '1' else '二'}学期"
        else:
            semester_display = str(raw_sem)
    else:
        semester_display = str(raw_sem)

    # 授课时间从 classes 提取
    classes = course.get("classes", [])
    date_str = ""
    if classes and isinstance(classes, list):
        c0 = classes[0]
        date_str = c0.get("schedule_time", "")

    info = {
        "course_name": course.get("name", course.get("course_name", "")),
        "course_code": course.get("code", course.get("course_code", "")),
        "teacher": teacher_name,
        "semester": semester_display,
        "date": date_str,
        "email": teacher_obj.get("email", "") if isinstance(teacher_obj, dict) else "",
        "office": teacher_office,
    }

    # 从 calendar 提取某周的信息
    if week is not None:
        calendar = data.get("calendar", [])
        week_data = None
        for entry in calendar:
            if entry.get("week") == week:
                week_data = entry
                break

        if week_data:
            info["week"] = week
            info["chapter_title"] = week_data.get("chapter_title", "")
            info["topic"] = week_data.get("topic", "")
            info["content"] = week_data.get("content", "")
            info["task"] = week_data.get("task", "")

            # 提取上一周的内容作为回顾
            prev_week = week - 1
            for entry in calendar:
                if entry.get("week") == prev_week:
                    info["prev_topic"] = entry.get("topic", "")
                    info["prev_content"] = entry.get("content", "")
                    break

    # 从 textbooks 提取引用
    textbooks = data.get("textbooks", [])
    if textbooks:
        info["references"] = [tb.get("citation", "") for tb in textbooks if tb.get("citation")]

    return info


# ====================================================================
# H2 标题提取（从裸 PPT 的 Module 过渡页中扫描）
# ====================================================================

def _extract_h2_from_script(pptx_path: str, info: dict) -> list:
    """从输入 PPTX 中提取 H2 过渡页的模块标题列表。

    generate_course_ppt.js 为每个 H2 跳变生成一张过渡页，
    特征：深色背景 + 单个大字号居中文本 + 无 Speaker Notes。
    我们也兼容直接从 compiled.md 提取。

    优先级：
    1. 从 package.yaml 同目录的 compiled.md 提取 H2 行（最可靠）
    2. 从 PPTX 的过渡页扫描（fallback）
    """
    import re

    # --- 策略 1：从 compiled.md 提取 ---
    # 推导 compiled.md 路径：
    # 输入 PPTX 路径类似 .../build/artifacts/_intermediate/课程_周次_Presentation.pptx
    # compiled.md 在 .../weeks/<weekId>/.build/compiled.md
    pptx_p = Path(pptx_path).resolve()

    # 从文件名推导周次 ID (e.g. 信息可视化_W01_Visual_Perception_Presentation.pptx)
    stem = pptx_p.stem
    # 匹配 <课程>_<W0X_Name>_Presentation
    m = re.match(r'.+?_(W\d+_.+?)_Presentation$', stem)
    if m:
        week_id = m.group(1)
        # 向上找课程根目录（build/artifacts/_intermediate/ 的 3 层父级）
        course_root = pptx_p.parent.parent.parent.parent
        compiled_path = course_root / "weeks" / week_id / ".build" / "compiled.md"
        if compiled_path.exists():
            h2_titles = []
            with open(compiled_path, "r", encoding="utf-8") as f:
                for line in f:
                    line = line.strip()
                    # 匹配 ## 开头但不匹配 ### 或 ####
                    if re.match(r'^##(?!#)\s+(.+)$', line):
                        title = re.match(r'^##\s+(.+)$', line).group(1).strip()
                        h2_titles.append(title)
            if h2_titles:
                return h2_titles

    # --- 策略 1b：从 src/ 目录扫描 M*.md 文件的 H2 行 ---
    if m:
        src_dir = course_root / "weeks" / week_id / "src"
        if src_dir.exists():
            h2_titles = []
            # 按文件名排序确保 M00, M01, M02... 顺序
            md_files = sorted(src_dir.glob("M*.md"))
            for md_file in md_files:
                with open(md_file, "r", encoding="utf-8") as f:
                    for line in f:
                        line = line.strip()
                        if re.match(r'^##(?!#)\s+(.+)$', line):
                            title = re.match(r'^##\s+(.+)$', line).group(1).strip()
                            h2_titles.append(title)
                            break  # 每个文件只取第一个 H2
            if h2_titles:
                return h2_titles

    # --- 策略 2：从 PPTX 过渡页扫描 ---
    try:
        from pptx import Presentation as PrsReader
        prs = PrsReader(pptx_path)
        h2_titles = []
        for slide in prs.slides:
            # 过渡页特征：无 notes，只有 1-2 个 shape，文本单行
            notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
            if notes.strip():
                continue  # 有 notes 的不是过渡页
            text_shapes = [s for s in slide.shapes if s.has_text_frame]
            if len(text_shapes) == 1:
                txt = text_shapes[0].text_frame.text.strip()
                # 过滤掉太短或太长的（过渡页标题通常 5~80 字符）
                if 5 <= len(txt) <= 80 and "\n" not in txt:
                    h2_titles.append(txt)
        if h2_titles:
            return h2_titles
    except Exception:
        pass

    return None


# ====================================================================
# CLI 入口
# ====================================================================

def main():
    ap = argparse.ArgumentParser(
        description="NFU PPT 品牌注入 — 固定环节包裹模式",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例：
  # 基础包裹
  %(prog)s -i slides.pptx -o out.pptx --course-name "交互设计" --teacher "张三"

  # 从 course.yaml + 周次自动填充（推荐）
  %(prog)s -i slides.pptx -o out.pptx --course-yaml course.yaml --week 3

  # 跳过回顾页和目录页
  %(prog)s -i slides.pptx -o out.pptx --skip review,toc --course-name ...

  # 添加章节分隔页
  %(prog)s -i slides.pptx -o out.pptx --sections "第1章：基础,第2章：进阶" ...

  # 仅替换主题色
  %(prog)s -i slides.pptx -o out.pptx --theme-only
""")

    ap.add_argument("-i", "--input", required=True, help="输入 .pptx")
    ap.add_argument("-o", "--output", required=True, help="输出 .pptx")

    # 课程信息
    g = ap.add_argument_group("课程信息")
    g.add_argument("--course-name", help="课程名称")
    g.add_argument("--course-code", help="课程编号")
    g.add_argument("--teacher", help="主讲教师")
    g.add_argument("--semester", help="授课学期")
    g.add_argument("--date", help="授课时间")
    g.add_argument("--email", help="教师邮箱")
    g.add_argument("--office", help="办公地点")
    g.add_argument("--course-yaml", help="从 course.yaml 读取")
    g.add_argument("--week", type=int, default=None,
                   help="周次编号，与 --course-yaml 配合使用，"
                        "自动填充回顾/目录/引用/作业内容")

    # 功能控制
    f = ap.add_argument_group("功能控制")
    f.add_argument("--theme-only", action="store_true",
                   help="仅注入主题色板，不增减幻灯片")
    f.add_argument("--no-theme", action="store_true",
                   help="不替换主题色板")
    f.add_argument("--skip", default="",
                   help="跳过的固定环节（逗号分隔）: review,toc,reference,assignment")
    f.add_argument("--sections", default="",
                   help="章节分隔页标题（逗号分隔），例如: '第1章：基础,第2章：进阶'")

    args = ap.parse_args()

    # 加载主题
    theme = load_theme()

    # 构建课程信息
    if args.course_yaml:
        info = load_course_yaml(args.course_yaml, week=args.week)
    else:
        info = {k: getattr(args, k) or "" for k in
                ["course_name", "course_code", "teacher",
                 "semester", "date", "email", "office"]}
    # 命令行覆盖 yaml
    for k in ["course_name", "course_code", "teacher"]:
        v = getattr(args, k)
        if v:
            info[k] = v
    if not info.get("course_name"):
        info["course_name"] = "课程名称"

    # 加载 PPT
    prs = Presentation(args.input)
    aspect = detect_aspect(prs)
    layout_cfg = theme["layout"][aspect]
    media_dir = resolve_media_dir(aspect)
    skip = set(s.strip() for s in args.skip.split(",") if s.strip())

    # W01 无上周数据 → 自动跳过回顾页
    if args.week == 1 and "review" not in skip:
        skip.add("review")
        print("ℹ️ W01 无上周数据，自动跳过回顾页")

    print(f"📐 比例: {aspect} | 现有页数: {len(prs.slides)}"
          + (f" | 周次: W{args.week:02d}" if args.week else ""))

    # ── 仅主题色模式 ─────────────────────────
    if args.theme_only:
        inject_theme_colors(prs, theme)
        prs.save(args.output)
        print(f"✅ 仅主题色注入完成 → {args.output}")
        return

    # ── 注入主题色 ──────────────────────────
    if not args.no_theme:
        try:
            inject_theme_colors(prs, theme)
            print("✅ 主题色板已替换")
        except Exception as e:
            print(f"⚠️ 主题色注入失败（非致命）: {e}")

    # ── 尾部固定环节（先添加，因为头部会移动索引）───
    original_count = len(prs.slides)

    if "reference" not in skip:
        refs = info.get("references")
        gen_reference(prs, theme, layout_cfg, media_dir, info,
                      references=refs)
        print(f"✅ 引用页" + (f" ({len(refs)} 条)" if refs else ""))

    if "assignment" not in skip:
        task_text = info.get("task")
        assignment_lines = None
        if task_text:
            assignment_lines = [task_text]
        gen_assignment(prs, theme, layout_cfg, media_dir, info,
                       assignment_text=assignment_lines)
        print("✅ 作业页" + (" (已填充)" if task_text else ""))

    gen_ending(prs, theme, layout_cfg, media_dir, info)
    print("✅ 结尾页")

    # ── 头部固定环节（从后往前插入到位置 0）────

    if "toc" not in skip:
        # 优先使用 --sections CLI，否则从实际脚本 H2 标题提取
        sections_list = None
        if args.sections:
            sections_list = [s.strip() for s in args.sections.split(",")]
        else:
            # 尝试从编译后的脚本中提取 H2 标题（Module 级目录）
            sections_list = _extract_h2_from_script(args.input, info)
            if not sections_list and info.get("content"):
                # 降级：从 course.yaml 的 calendar content 提取
                raw = info["content"].strip()
                sections_list = [line.strip() for line in raw.split("\n")
                                 if line.strip()]
        gen_toc(prs, theme, layout_cfg, media_dir, info, sections=sections_list)
        _move_slide(prs, len(prs.slides) - 1, 0)
        print("✅ 目录页" + (f" ({len(sections_list)} 条)" if sections_list else ""))

    if "review" not in skip:
        # 从 calendar 前一周提取回顾内容
        review_title = "上节课知识点回顾"
        review_points = None
        if info.get("prev_topic"):
            review_title = f"回顾：{info['prev_topic']}"
            if info.get("prev_content"):
                review_points = [line.strip()
                                 for line in info["prev_content"].strip().split("\n")
                                 if line.strip()]
        gen_review(prs, theme, layout_cfg, media_dir, info,
                   title=review_title, points=review_points)
        _move_slide(prs, len(prs.slides) - 1, 0)
        print(f"✅ 回顾页：{review_title}")

    # 封面页 — 始终生成
    gen_cover(prs, theme, layout_cfg, media_dir, info, position=0)
    print("✅ 封面页")

    # ── 可选：章节分隔页（仅 CLI --sections 时）──
    if args.sections:
        section_titles = [s.strip() for s in args.sections.split(",")]
        header_count = 1  # cover
        if "review" not in skip:
            header_count += 1
        if "toc" not in skip:
            header_count += 1
        for idx, st in enumerate(section_titles):
            gen_section_header(prs, theme, layout_cfg, media_dir,
                               chapter=st, section_title="")
            _move_slide(prs, len(prs.slides) - 1, header_count + idx)
        print(f"✅ {len(section_titles)} 个章节分隔页")

    # ── 保存 ──────────────────────────────
    prs.save(args.output)
    print(f"\n🎉 品牌化完成 → {args.output}")
    print(f"   总页数: {len(prs.slides)} "
          f"(原有 {original_count} + 固定环节 {len(prs.slides) - original_count})")


if __name__ == "__main__":
    main()
