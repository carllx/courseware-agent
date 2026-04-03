#!/usr/bin/env python3
"""
sync_notes_back.py — PPT Presentation Notes → Markdown 脚本逆向同步

将 PPT 中修改过的演讲备注（Presentation Notes）同步回 Markdown 脚本。
基于 Slide ID 锚点进行精确映射。

用法:
    python engines/sync_notes_back.py \
        --pptx "信息可视化/build/presentations/W01_Presentation_Gen.pptx" \
        --script "信息可视化/scripts/W01_Visual_Perception.md" \
        --dry-run

    # 实际覆盖，并创建备份 (.bak)
    python engines/sync_notes_back.py \
        --pptx ... --script ... \
        --extra-tags "NOTE" "REF" "HINT"

依赖: python-pptx
"""

import argparse
import os
import re
import shutil
import sys
import difflib
import logging
from datetime import datetime

# 将 validation_suite 的 script_parser 加入路径
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
WORKSPACE_ROOT = os.path.dirname(SCRIPT_DIR)
PARSER_DIR = os.path.join(WORKSPACE_ROOT, '.agent', 'skills', 'validation_suite', 'scripts')
sys.path.insert(0, PARSER_DIR)

from script_parser import parse_script, BlockType

# ============================================================
# 日志配置
# ============================================================
logger = logging.getLogger('sync_notes_back')

# ============================================================
# 默认保护标签列表（可通过 --extra-tags 扩展）
# ============================================================
DEFAULT_PROTECTED_TAGS = [
    'TECH NOTE', 'WARNING', 'DID YOU KNOW', 'STORY TIME', 'PHILOSOPHY',
    'CASE STUDY', 'LIFE CONNECT', 'TEACHING MOMENT', 'ACTIVITY',
    'STAGE NOTE', '!INFO', '!TIP', '!WARNING', '!CAUTION', '!IMPORTANT', '!NOTE',
]


def build_protected_regex(extra_tags: list[str] | None = None) -> re.Pattern:
    """
    构建保护块起始正则。将默认标签与用户自定义标签合并。
    """
    tags = list(DEFAULT_PROTECTED_TAGS)
    if extra_tags:
        tags.extend(extra_tags)
    # 转义特殊字符并用 | 连接
    escaped = [re.escape(t) for t in tags]
    pattern = r'^\>\s*\[(' + '|'.join(escaped) + r')'
    return re.compile(pattern, re.IGNORECASE)


def extract_pptx_notes(pptx_path: str) -> list[str]:
    """
    从 PPTX 中提取每一页的 Presentation Notes 文本。
    返回按页序排列的文本列表。
    """
    from pptx import Presentation

    prs = Presentation(pptx_path)
    notes = []
    for slide in prs.slides:
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            text = slide.notes_slide.notes_text_frame.text.strip()
            notes.append(text)
        else:
            notes.append("")
    return notes


def build_speech_map(
    script_path: str,
    extra_tags: list[str] | None = None,
) -> list[dict]:
    """
    解析 Markdown 脚本，为每个 VISUAL 块建立 Speech 区域映射。
    将 Speech 区域拆分为「纯 Speech 行」和「保护块」(标签/ACTIVITY 等)。

    返回列表，每项包含:
    {
        'slide_id': str,
        'speech_line_start': int,       # 整体区域起始行 (1-indexed)
        'speech_line_end': int,         # 整体区域结束行 (1-indexed, inclusive)
        'speech_text': str,             # 纯 Speech 文本（用于比较）
        'speech_line_numbers': [int],   # 纯 Speech 行号列表 (1-indexed)
        'protected_blocks': [(int,int,[str])],  # 保护块: (start,end,lines) 列表
    }
    """
    blocks = parse_script(script_path)
    result = []

    with open(script_path, 'r', encoding='utf-8') as f:
        all_lines = f.readlines()

    visual_indices = [i for i, b in enumerate(blocks) if b.block_type == BlockType.VISUAL]
    re_non_speech = build_protected_regex(extra_tags)

    for vi, block_idx in enumerate(visual_indices):
        vblock = blocks[block_idx]
        slide_id = vblock.metadata.get('slide_id', f'(无ID-slide{vi + 1})')

        # 确定 Speech 区域的行范围
        speech_start = vblock.line_end + 1

        if vi + 1 < len(visual_indices):
            next_visual_idx = visual_indices[vi + 1]
            speech_end = blocks[next_visual_idx].line_start - 1
        else:
            speech_end = len(all_lines)
            for line_idx in range(vblock.line_end, len(all_lines)):
                if all_lines[line_idx].strip().startswith('## Self-Verification'):
                    speech_end = line_idx
                    break

        # 遍历范围，分类每一行为 Speech 或 Protected
        speech_lines = []           # 纯 Speech 文本
        speech_line_numbers = []    # 纯 Speech 行号
        protected_blocks = []       # 保护块列表 [(start, end, [raw_lines])]

        in_protected = False
        protected_start = None
        protected_lines_buf = []

        for line_num in range(speech_start, speech_end + 1):
            if line_num > len(all_lines):
                break
            raw = all_lines[line_num - 1]
            stripped = raw.strip()

            # 检测保护块开始
            if re_non_speech.match(stripped):
                # 如果之前有未关闭的保护块，先关闭
                if in_protected and protected_lines_buf:
                    protected_blocks.append((protected_start, line_num - 1, list(protected_lines_buf)))
                in_protected = True
                protected_start = line_num
                protected_lines_buf = [raw]
                continue

            # 在保护块内 —— 任何以 > 开头的行、空行、
            # 以及结构标记（---、#、<!-- -->）都视为块的一部分
            if in_protected:
                if stripped.startswith('>') or stripped == '':
                    protected_lines_buf.append(raw)
                    continue
                elif stripped.startswith('---') or stripped.startswith('#') or stripped.startswith('<!--'):
                    # 结构标记在保护块内部 → 仍属于保护块
                    protected_lines_buf.append(raw)
                    continue
                else:
                    # 保护块结束（当前行是非引用/非结构行）
                    protected_blocks.append((protected_start, line_num - 1, list(protected_lines_buf)))
                    in_protected = False
                    protected_start = None
                    protected_lines_buf = []
                    # 当前行继续按 Speech 逻辑处理（fall through）

            # 跳过分隔线、标题、注释等（不视为 Speech 也不保护）
            if stripped.startswith('---') or stripped.startswith('#'):
                continue
            if stripped.startswith('<!--'):
                continue
            if stripped == '[SPEECH]':
                continue

            # 纯 Speech 行
            speech_lines.append(raw.rstrip('\n'))
            speech_line_numbers.append(line_num)

        # 关闭末尾未关闭的保护块
        if in_protected and protected_lines_buf:
            protected_blocks.append((protected_start, speech_end, list(protected_lines_buf)))

        speech_text = '\n'.join(speech_lines).strip()

        result.append({
            'slide_id': slide_id,
            'visual_line_start': vblock.line_start,
            'speech_line_start': speech_start,
            'speech_line_end': speech_end,
            'speech_text': speech_text,
            'speech_line_numbers': speech_line_numbers,
            'protected_blocks': protected_blocks,
        })

    return result


def normalize_for_compare(text: str) -> str:
    """
    标准化文本用于比较。
    移除 Markdown 标记、多余空格，使 PPT Notes 和原始 Speech 可比。
    """
    # 全角空格 → 半角
    text = text.replace('\u3000', ' ')
    # Markdown 格式标记
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    text = re.sub(r'`(.*?)`', r'\1', text)
    text = re.sub(r'\(Pause:.*?\)', '', text)
    text = re.sub(r'^\s*[\-\*]\s+', '• ', text, flags=re.MULTILINE)
    # 合并连续空行
    text = re.sub(r'\n{3,}', '\n\n', text)
    # 合并空白
    text = re.sub(r'[ \t]+', ' ', text)
    return text.strip()


def clean_new_text(text: str) -> str:
    """
    清理即将写回的新 Speech 文本：
    - 全角空格 → 半角
    - 合并连续空行（最多保留 1 个）
    - 统一换行符为 \n
    """
    # 统一换行符
    text = text.replace('\r\n', '\n').replace('\r', '\n')
    # 全角空格 → 半角
    text = text.replace('\u3000', ' ')
    # 合并连续空行
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def create_timestamped_backup(script_path: str) -> str:
    """
    创建带时间戳的备份文件，避免覆盖历史备份。
    返回备份文件路径。
    """
    timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
    base, ext = os.path.splitext(script_path)
    bak_path = f'{base}.bak_{timestamp}{ext}'
    shutil.copy2(script_path, bak_path)
    return bak_path


def main():
    parser = argparse.ArgumentParser(
        description="PPT Presentation Notes → Markdown 脚本逆向同步"
    )
    parser.add_argument('--pptx', required=True, help="PPT 文件路径")
    parser.add_argument('--script', required=True, help="Markdown 脚本文件路径")
    parser.add_argument('--dry-run', action='store_true',
                        help="仅输出差异报告，不修改文件")
    parser.add_argument('--extra-tags', nargs='*', default=[],
                        help="额外需要保护的标签名（如 NOTE REF HINT）")
    parser.add_argument('-v', '--verbose', action='store_true',
                        help="启用详细日志输出")
    args = parser.parse_args()

    # 配置日志
    log_level = logging.DEBUG if args.verbose else logging.WARNING
    logging.basicConfig(level=log_level, format='%(levelname)s: %(message)s')

    pptx_path = os.path.abspath(args.pptx)
    script_path = os.path.abspath(args.script)

    if not os.path.exists(pptx_path):
        print(f"❌ PPT 文件不存在: {pptx_path}")
        sys.exit(1)
    if not os.path.exists(script_path):
        print(f"❌ 脚本文件不存在: {script_path}")
        sys.exit(1)

    print(f"📖 PPT: {os.path.basename(pptx_path)}")
    print(f"📝 脚本: {os.path.basename(script_path)}")
    if args.extra_tags:
        print(f"🏷️  额外保护标签: {', '.join(args.extra_tags)}")
    print(f"{'🔍 DRY RUN 模式' if args.dry_run else '✏️  实际同步模式'}")
    print()

    # 1. 提取 PPT Notes（带异常捕获）
    try:
        pptx_notes = extract_pptx_notes(pptx_path)
    except Exception as e:
        print(f"❌ 读取 PPT 文件失败: {e}")
        logger.debug("PPT 读取异常详情:", exc_info=True)
        sys.exit(1)
    print(f"📊 PPT Slide 数: {len(pptx_notes)}")

    # 2. 解析脚本 Speech 映射（带异常捕获）
    try:
        speech_map = build_speech_map(script_path, extra_tags=args.extra_tags or None)
    except Exception as e:
        print(f"❌ 解析脚本失败: {e}")
        logger.debug("脚本解析异常详情:", exc_info=True)
        sys.exit(1)
    print(f"📊 脚本 VISUAL 块数: {len(speech_map)}")

    # 3. 数量校验
    if len(pptx_notes) != len(speech_map):
        print(f"\n⚠️  数量不匹配！PPT 有 {len(pptx_notes)} 页，脚本有 {len(speech_map)} 个 VISUAL 块")
        print("   可能原因: PPT 中增删了幻灯片，或脚本中有未被识别的 VISUAL 块")
        print("   中止同步，请手动对齐后重试")
        sys.exit(1)

    # 4. 逐 Slide 比较差异
    changes = []
    for i, (note, mapping) in enumerate(zip(pptx_notes, speech_map)):
        slide_id = mapping['slide_id']
        original = mapping['speech_text']

        norm_note = normalize_for_compare(note)
        norm_original = normalize_for_compare(original)

        if norm_note == norm_original:
            continue
        if not note.strip():
            continue

        diff = list(difflib.unified_diff(
            original.splitlines(keepends=True),
            note.splitlines(keepends=True),
            fromfile=f'脚本/{slide_id}',
            tofile=f'PPT/{slide_id}',
            lineterm=''
        ))

        changes.append({
            'index': i,
            'slide_id': slide_id,
            'line_start': mapping['speech_line_start'],
            'line_end': mapping['speech_line_end'],
            'speech_line_numbers': mapping['speech_line_numbers'],
            'protected_blocks': mapping['protected_blocks'],
            'original': original,
            'new_text': clean_new_text(note),
            'diff': diff,
        })

    # 5. 报告
    print(f"\n{'=' * 60}")
    if not changes:
        print("✅ 无差异 — PPT Notes 与脚本 Speech 完全一致")
        sys.exit(0)

    print(f"📝 发现 {len(changes)} 处差异:\n")
    for c in changes:
        n_protected = len(c['protected_blocks'])
        prot_info = f"  🛡️  {n_protected} 个保护块将被保留" if n_protected else ""
        print(f"  Slide {c['index'] + 1} [{c['slide_id']}] (脚本行 {c['line_start']}-{c['line_end']}){prot_info}")
        for dl in c['diff'][:10]:
            print(f"    {dl}")
        if len(c['diff']) > 10:
            print(f"    ... ({len(c['diff']) - 10} 行省略)")
        print()

    if args.dry_run:
        print(f"🔍 DRY RUN 完成 — 未修改任何文件")
        sys.exit(0)

    # 6. 执行安全同步（带异常捕获）
    try:
        bak_path = create_timestamped_backup(script_path)
        print(f"💾 备份已创建: {os.path.basename(bak_path)}")
    except Exception as e:
        print(f"❌ 创建备份失败: {e}")
        sys.exit(1)

    try:
        with open(script_path, 'r', encoding='utf-8') as f:
            lines = f.readlines()

        # 检测原始换行符风格
        line_ending = '\n'
        if lines and lines[0].endswith('\r\n'):
            line_ending = '\r\n'

        # 从后向前处理（避免行号偏移）
        for c in reversed(changes):
            region_start = c['line_start'] - 1   # 0-indexed inclusive
            region_end = c['line_end']            # 0-indexed exclusive

            # 收集所有保护块的行号集合
            protected_line_set = set()
            for (pb_start, pb_end, pb_lines) in c['protected_blocks']:
                for ln in range(pb_start, pb_end + 1):
                    protected_line_set.add(ln)

            # 重建区域：遍历原始行，保护块原样保留，Speech 行替换
            new_region = []
            speech_inserted = False

            for line_num in range(c['line_start'], c['line_end'] + 1):
                if line_num > len(lines):
                    break

                raw = lines[line_num - 1]
                stripped = raw.strip()

                # 保护块行 → 原样保留
                if line_num in protected_line_set:
                    # 在保护块之前插入新 Speech（仅一次）
                    if not speech_inserted:
                        new_region.append(line_ending)
                        for sl in c['new_text'].split('\n'):
                            new_region.append(sl + line_ending)
                        new_region.append(line_ending)
                        speech_inserted = True
                    new_region.append(raw)
                    continue

                # 分隔线 / 标题 / 注释 → 原样保留
                if stripped.startswith('---') or stripped.startswith('#') or stripped.startswith('<!--'):
                    # 在结构标记之前插入新 Speech（仅一次）
                    if not speech_inserted:
                        new_region.append(line_ending)
                        for sl in c['new_text'].split('\n'):
                            new_region.append(sl + line_ending)
                        new_region.append(line_ending)
                        speech_inserted = True
                    new_region.append(raw)
                    continue

                # [SPEECH] 标记 → 保留
                if stripped == '[SPEECH]':
                    new_region.append(raw)
                    continue

                # 纯 Speech 行 → 跳过（将被新文本替换）

            # 如果 Speech 尚未插入（区域内全部是纯 Speech 行，无保护块）
            if not speech_inserted:
                new_region.append(line_ending)
                for sl in c['new_text'].split('\n'):
                    new_region.append(sl + line_ending)
                new_region.append(line_ending)

            # 替换区域
            lines[region_start:region_end] = new_region

        with open(script_path, 'w', encoding='utf-8') as f:
            f.writelines(lines)

    except Exception as e:
        print(f"❌ 写入文件失败: {e}")
        print(f"   原始文件已备份于: {os.path.basename(bak_path)}")
        logger.debug("写入异常详情:", exc_info=True)
        sys.exit(1)

    total_protected = sum(len(c['protected_blocks']) for c in changes)
    print(f"✅ 同步完成 — 已更新 {len(changes)} 处 Speech 内容")
    print(f"   🛡️  保护了 {total_protected} 个结构化块（TECH NOTE / ACTIVITY 等）")
    print(f"   原始文件备份于: {os.path.basename(bak_path)}")


if __name__ == '__main__':
    main()
